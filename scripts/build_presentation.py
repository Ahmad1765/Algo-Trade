"""
Build the Algo-Trade class presentation (.pptx) for a Simulation & Modeling course.

Generates a polished, dark "trading terminal" themed deck using python-pptx.
Run:  python scripts/build_presentation.py
Output:  docs/Algo-Trade-Simulation-and-Modeling.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
BG        = RGBColor(0x0A, 0x0A, 0x0B)   # near-black (zinc-950)
PANEL     = RGBColor(0x16, 0x16, 0x19)   # card panel
PANEL2    = RGBColor(0x1F, 0x1F, 0x24)   # lighter panel
STROKE    = RGBColor(0x2A, 0x2A, 0x31)   # subtle border
WHITE     = RGBColor(0xF4, 0xF4, 0xF5)   # zinc-100
MUTED     = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400
DIM       = RGBColor(0x71, 0x71, 0x7A)   # zinc-500
EMERALD   = RGBColor(0x10, 0xB9, 0x81)   # accent green
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # accent amber
RED       = RGBColor(0xEF, 0x44, 0x44)   # danger
BLUE      = RGBColor(0x38, 0xBD, 0xF8)   # info

FONT      = "Segoe UI"
MONO      = "Consolas"

# 16:9 widescreen
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _noline(shape):
    shape.line.fill.background()


def _line(shape, color, w=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _solid(bg, BG)
    _noline(bg)
    bg.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        _solid(shp, fill)
    if line is None:
        _noline(shp)
    else:
        _line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            r.font.italic = italic
    return tb


def R(txt, size=18, color=WHITE, bold=False, font=FONT, italic=False):
    return (txt, size, color, bold, font, italic)


def accent_bar(s, x=0.0, y=0.0, w=0.18, h=7.5, color=EMERALD):
    rect(s, x, y, w, h, fill=color)


def kicker(s, label, x=0.7, y=0.55, color=EMERALD):
    text(s, x, y, 11, 0.4, [[R(label.upper(), 13, color, True, MONO)]])


def title(s, t, x=0.7, y=0.95, w=12, size=40, color=WHITE):
    text(s, x, y, w, 1.2, [[R(t, size, color, True)]])


def footer(s, idx, total=14, tag="ALGO-TRADE · Simulation & Modeling"):
    text(s, 0.7, 7.02, 9, 0.35, [[R(tag, 10, DIM, False, MONO)]])
    text(s, 11.4, 7.02, 1.2, 0.35, [[R(f"{idx:02d} / {total:02d}", 10, DIM, False, MONO)]],
         align=PP_ALIGN.RIGHT)


def chip(s, x, y, label, color=EMERALD, w=2.0):
    c = rect(s, x, y, w, 0.42, fill=PANEL2, line=color, line_w=1.0, radius=True)
    text(s, x, y + 0.045, w, 0.33,
         [[R(label, 12, color, True, MONO)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def card(s, x, y, w, h, head, head_color, body_lines, body_size=13):
    rect(s, x, y, w, h, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
    rect(s, x, y, 0.07, h, fill=head_color)
    text(s, x + 0.28, y + 0.22, w - 0.5, 0.5, [[R(head, 16, WHITE, True)]])
    paras = [[R(line, body_size, MUTED)] for line in body_lines]
    text(s, x + 0.28, y + 0.78, w - 0.5, h - 1.0, paras, space_after=5,
         line_spacing=1.05)


TOTAL = 14
n = 0


def nxt():
    global n
    n += 1
    return n


# ----------------------------------------------------------------------------
# Slide 1 — Title
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=EMERALD)
# faint grid panel right
rect(s, 8.4, 0, 4.93, 7.5, fill=PANEL)
rect(s, 8.4, 0, 0.02, 7.5, fill=STROKE)

text(s, 0.9, 1.7, 8, 0.5, [[R("SIMULATION & MODELING  ·  FINAL PROJECT", 14, EMERALD, True, MONO)]])
text(s, 0.9, 2.35, 8.2, 2.2, [
    [R("Algo-Trade", 64, WHITE, True)],
    [R("Modeling & Simulating the", 30, MUTED, False)],
    [R("Options Market", 30, WHITE, True)],
])
text(s, 0.9, 4.95, 7.3, 1.2, [[R(
    "An event-driven system that models market behavior with mathematical "
    "indicators and simulates a full trading strategy end-to-end — without "
    "risking a single real dollar.", 15, MUTED)]], line_spacing=1.15)

text(s, 0.9, 6.35, 8, 1.0, [
    [R("Presented by   ", 13, DIM, False, MONO),
     R("Course: Simulation and Modeling", 13, EMERALD, True, MONO)],
    [R("Muhammad Ahmad Saleem  F2022065202", 12.5, WHITE, False, MONO)],
    [R("Danish Salman  F2022065114      Irha Shoaib  F2021065136", 12.5, WHITE, False, MONO)],
], space_after=3, line_spacing=1.0)

# right-side "terminal" mock
text(s, 8.75, 0.7, 4.4, 0.4, [[R("● ● ●   live-paper-trade", 12, DIM, False, MONO)]])
term = [
    [R("$ python -m src.cli.main --mode paper", 12, EMERALD, False, MONO)],
    [R("[boot] event bus online", 11, MUTED, False, MONO)],
    [R("[scan] top movers: 10 symbols", 11, MUTED, False, MONO)],
    [R("[model] RSI=72.4  MACD=+0.31", 11, BLUE, False, MONO)],
    [R("[signal] AAPL CALL  conf=0.81", 11, AMBER, False, MONO)],
    [R("[risk]  size=4.0% equity  OK", 11, MUTED, False, MONO)],
    [R("[order] FILLED @ 2.14  (sim)", 11, EMERALD, False, MONO)],
    [R("[cb]    daily P&L +1.2%  armed", 11, MUTED, False, MONO)],
]
text(s, 8.75, 1.4, 4.4, 5, term, space_after=8, line_spacing=1.0)
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 2 — The Problem / Motivation
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=AMBER)
kicker(s, "Motivation", color=AMBER)
title(s, "Why model a market instead of guessing?")

text(s, 0.7, 2.0, 7.1, 3.5, [
    [R("Human trading is slow, emotional, and inconsistent.", 18, WHITE, True)],
    [R("", 8, MUTED)],
    [R("A stock/option market is a complex, noisy, fast-moving system. "
       "People react to fear and greed; they can't watch hundreds of "
       "symbols at once or apply the same rules every time.", 15, MUTED)],
    [R("", 8, MUTED)],
    [R("The Simulation & Modeling question:", 16, EMERALD, True)],
    [R("Can we capture trading decisions as a precise model, then "
       "simulate that model against the market to test it safely — "
       "before any real money is on the line?", 15, MUTED)],
], line_spacing=1.15)

card(s, 8.1, 2.0, 4.6, 1.55, "The model", BLUE,
     ["Turn price action into math:", "indicators + rules → a signal."])
card(s, 8.1, 3.7, 4.6, 1.55, "The simulation", EMERALD,
     ["Replay history & run live paper", "trades with a mock broker."])
card(s, 8.1, 5.4, 4.6, 1.05, "The safety net", RED,
     ["Risk limits + circuit breaker", "stop runaway behavior."])
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 3 — What is Algo-Trade
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "Overview")
title(s, "What is Algo-Trade?")
text(s, 0.7, 1.95, 11.9, 0.9, [[R(
    "A production-grade, event-driven algorithmic options-trading system written in "
    "Python — built strictly for education and paper trading.", 17, MUTED)]],
    line_spacing=1.15)

feats = [
    ("Scans the market", "Pulls top gainers/losers in real time", EMERALD),
    ("Filters options", "Volume, open interest, spread, expiry", BLUE),
    ("Generates signals", "RSI + MACD momentum confirmation", AMBER),
    ("Plans the trade", "Entry, stop-loss, take-profit via ATR", EMERALD),
    ("Manages risk", "Position sizing, PDT, daily breaker", RED),
    ("Simulates / executes", "Mock broker, backtest, or live", BLUE),
]
x0, y0, cw, ch, gx, gy = 0.7, 3.0, 3.86, 1.55, 0.18, 0.2
for i, (h_, b_, c_) in enumerate(feats):
    r, col = divmod(i, 3)
    card(s, x0 + col * (cw + gx), y0 + r * (ch + gy), cw, ch, h_, c_, [b_], body_size=13)
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 4 — Mapping to Simulation & Modeling concepts
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=BLUE)
kicker(s, "Course Connection", color=BLUE)
title(s, "How it maps to Simulation & Modeling")

rows = [
    ("Mathematical model", "RSI, MACD, ATR turn raw prices into bounded, computable indicators", EMERALD),
    ("Discrete-event simulation", "An event bus drives the system: tick → signal → order → fill", BLUE),
    ("System under test", "Strategy rules are the model we evaluate for performance", AMBER),
    ("Historical simulation", "Backtester replays recorded minute-bars to measure outcomes", EMERALD),
    ("Real-time simulation", "Paper mode runs live data through a mock broker — no real fills", BLUE),
    ("Stochastic input", "Mock market adapter generates synthetic price paths for testing", AMBER),
    ("Validation & metrics", "Win rate, P&L, signal counts quantify model quality", EMERALD),
]
y = 2.05
for label, desc, c in rows:
    rect(s, 0.7, y, 0.12, 0.62, fill=c)
    text(s, 1.0, y + 0.02, 3.7, 0.6, [[R(label, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.8, y + 0.02, 7.8, 0.6, [[R(desc, 14, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.71
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 5 — Architecture pipeline
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "Architecture")
title(s, "Event-driven pipeline")
text(s, 0.7, 1.9, 12, 0.5, [[R(
    "Each stage emits a typed event onto the bus; the next stage reacts. Components "
    "are decoupled and independently testable.", 14, MUTED)]])

stages = [
    ("Market\nData", "FMP · Yahoo · Mock", BLUE),
    ("Screener", "Top gainers/losers", EMERALD),
    ("Options\nFetcher", "Liquidity filter", BLUE),
    ("Indicators", "RSI · MACD · ATR", AMBER),
    ("Strategy\nEngine", "Signal + confirm", EMERALD),
    ("Risk\nManager", "Size · limits", RED),
    ("Execution", "Mock / Webull", BLUE),
]
bw, bh, gap = 1.58, 1.5, 0.22
x = 0.7
y = 2.95
for i, (name, sub, c) in enumerate(stages):
    rect(s, x, y, bw, bh, fill=PANEL, line=c, line_w=1.25, radius=True)
    text(s, x, y + 0.22, bw, 0.8, [[R(p, 14, WHITE, True)] for p in name.split("\n")],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x, y + bh - 0.5, bw, 0.4, [[R(sub, 9.5, MUTED, False, MONO)]],
         align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.01),
                                Inches(y + bh / 2 - 0.1), Inches(gap - 0.02), Inches(0.2))
        _solid(ar, DIM)
        _noline(ar)
        ar.shadow.inherit = False
    x += bw + gap

# bus line + side services
rect(s, 0.7, 4.85, 11.92, 0.5, fill=PANEL2, line=STROKE, line_w=1.0, radius=True)
text(s, 0.7, 4.9, 11.92, 0.4, [[R("◀  EVENT BUS  (typed events: Tick · Signal · Order · Fill)  ▶",
     13, MUTED, True, MONO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

card(s, 0.7, 5.6, 3.8, 1.1, "Persistence", BLUE, ["SQLite/SQLAlchemy — signals,", "positions, history"])
card(s, 4.76, 5.6, 3.8, 1.1, "API + Dashboard", EMERALD, ["aiohttp REST + Next.js UI", "with live SSE telemetry"])
card(s, 8.82, 5.6, 3.8, 1.1, "Observability", AMBER, ["Structured JSON logs", "with secret redaction"])
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 6 — The modeling layer (indicators)
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=AMBER)
kicker(s, "The Model", color=AMBER)
title(s, "Turning price into math: the indicators")

card(s, 0.7, 2.05, 3.86, 4.1, "RSI", EMERALD,
     ["Relative Strength Index", "", "Momentum oscillator, 0–100.", "",
      ">70 → overbought (CALL bias)", "<30 → oversold (PUT bias)", "",
      "Models whether a move is", "stretched and likely to turn."])
card(s, 4.74, 2.05, 3.86, 4.1, "MACD", BLUE,
     ["Moving Avg Convergence", "Divergence", "",
      "Difference of two EMAs vs a", "signal line.", "",
      "Crossovers confirm trend", "direction & strength —", "the confirmation filter."])
card(s, 8.78, 2.05, 3.86, 4.1, "ATR", AMBER,
     ["Average True Range", "", "Measures recent volatility.", "",
      "Used to size the trade plan:", "stop-loss & take-profit", "scale with how 'wide' the",
      "market is moving.", ""])

text(s, 0.7, 6.4, 12, 0.6, [[
    R("Pure functions → ", 14, MUTED), R("deterministic, unit-tested, reproducible", 14, EMERALD, True),
    R("  — the foundation of a trustworthy model.", 14, MUTED)]])
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 7 — The simulation engine
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "The Simulation")
title(s, "Four ways to run the model")

modes = [
    ("BACKTEST", EMERALD, "Historical simulation",
     "Replay recorded minute-bar CSVs through the full pipeline and print a P&L report. "
     "Fast, repeatable evaluation of the strategy."),
    ("PAPER", BLUE, "Real-time simulation",
     "Live (or synthetic) data flows through the system; the MockBrokerAdapter simulates "
     "fills. Full behavior, zero real money."),
    ("MANUAL", AMBER, "Decision support",
     "Signals are logged as recommendations only — no orders placed. The human stays in "
     "the loop."),
    ("AUTOMATED", RED, "Live execution",
     "Real orders via a broker adapter (Webull). Disabled for the demo — shown only to "
     "explain the deployment path."),
]
x0, y0, cw, ch, gx, gy = 0.7, 2.05, 5.86, 1.95, 0.2, 0.22
for i, (tag, c, sub, body) in enumerate(modes):
    r, col = divmod(i, 2)
    x, y = x0 + col * (cw + gx), y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
    rect(s, x, y, 0.07, ch, fill=c)
    chip(s, x + 0.28, y + 0.25, tag, color=c, w=1.7)
    text(s, x + 2.15, y + 0.27, cw - 2.4, 0.4, [[R(sub, 14, WHITE, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.28, y + 0.85, cw - 0.55, 1.0, [[R(body, 12.5, MUTED)]], line_spacing=1.1)
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 8 — Risk modeling
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=RED)
kicker(s, "Risk Modeling", color=RED)
title(s, "Modeling what can go wrong")
text(s, 0.7, 1.95, 12, 0.6, [[R(
    "A trading model is only useful if it can't blow up. Risk rules are first-class "
    "components in the simulation.", 15, MUTED)]], line_spacing=1.1)

card(s, 0.7, 2.95, 3.86, 3.3, "Position sizing", EMERALD,
     ["Each trade is capped at a", "% of equity (default 5%).", "",
      "No single bad bet can", "sink the account.", "",
      "max_position_pct"])
card(s, 4.74, 2.95, 3.86, 3.3, "Daily circuit breaker", AMBER,
     ["Halts all trading when:", "", "• profit target hit, or", "• daily loss limit hit", "",
      "Models a 'walk away'", "discipline rule."])
card(s, 8.78, 2.95, 3.86, 3.3, "PDT & limits", BLUE,
     ["Pattern-Day-Trader rule", "monitoring + guardrails.", "",
      "Keeps the simulated", "account within realistic", "regulatory constraints."])
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 9 — The UI / dashboard
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "The Interface")
title(s, "A live dashboard for the simulation")

# browser frame mock
fx, fy, fw, fh = 7.0, 2.0, 5.7, 4.3
rect(s, fx, fy, fw, fh, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
rect(s, fx, fy, fw, 0.45, fill=PANEL2)
for i, c in enumerate((RED, AMBER, EMERALD)):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(fx + 0.2 + i * 0.22), Inches(fy + 0.15),
                             Inches(0.13), Inches(0.13))
    _solid(dot, c); _noline(dot); dot.shadow.inherit = False
# stat tiles
tiles = [("MARKET", "OPEN", EMERALD), ("POSITIONS", "3", WHITE),
         ("SIGNALS", "27", AMBER), ("UPTIME", "04:12", BLUE)]
tw = (fw - 0.6 - 0.3) / 4
for i, (lab, val, c) in enumerate(tiles):
    tx = fx + 0.3 + i * (tw + 0.1)
    rect(s, tx, fy + 0.65, tw, 0.95, fill=PANEL2, line=STROKE, line_w=0.75, radius=True)
    text(s, tx, fy + 0.75, tw, 0.3, [[R(lab, 8, DIM, True, MONO)]], align=PP_ALIGN.CENTER)
    text(s, tx, fy + 1.05, tw, 0.4, [[R(val, 17, c, True)]], align=PP_ALIGN.CENTER)
# chart area
rect(s, fx + 0.3, fy + 1.75, fw - 0.6, 2.25, fill=PANEL2, line=STROKE, line_w=0.75, radius=True)
text(s, fx + 0.45, fy + 1.85, 3, 0.3, [[R("PRICE  ·  AAPL", 9, MUTED, True, MONO)]])
# simple zig-zag "chart" line using connectors
pts = [(0.3, 1.4), (1.0, 0.9), (1.6, 1.6), (2.3, 0.7), (3.1, 1.2),
       (3.8, 0.5), (4.6, 0.95), (5.1, 0.4)]
base_x = fx + 0.3
base_y = fy + 1.75 + 2.25 - 0.35
for i in range(len(pts) - 1):
    x1 = base_x + pts[i][0]
    y1 = base_y - pts[i][1]
    x2 = base_x + pts[i + 1][0]
    y2 = base_y - pts[i + 1][1]
    conn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = EMERALD
    conn.line.width = Pt(2)
    conn.shadow.inherit = False

# left bullet list
bullets = [
    ("Live dashboard", "stats, circuit-breaker banner, price chart, order panel"),
    ("Positions & Signals", "real-time tables of what the model is doing"),
    ("Backtest & Strategies", "configure and review simulation runs"),
    ("History & Settings", "audit trail + tune the model's parameters"),
]
y = 2.15
for h_, b_ in bullets:
    rect(s, 0.7, y + 0.06, 0.12, 0.5, fill=EMERALD)
    text(s, 1.0, y, 5.6, 0.4, [[R(h_, 16, WHITE, True)]])
    text(s, 1.0, y + 0.4, 5.7, 0.5, [[R(b_, 13, MUTED)]], line_spacing=1.0)
    y += 1.02
text(s, 0.7, 6.35, 5.9, 0.5, [[R("Next.js · React · Tailwind CSS · lucide icons",
     12, DIM, False, MONO)]])
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 10 — Tech stack
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=BLUE)
kicker(s, "Implementation", color=BLUE)
title(s, "Technology stack")

groups = [
    ("Backend", EMERALD, ["Python 3.11+", "asyncio event loop", "aiohttp REST API",
                          "SQLAlchemy + SQLite", "pytest test suite"]),
    ("Modeling", AMBER, ["RSI / MACD / ATR", "Pure-function indicators", "ATR-based trade plans",
                         "Config-driven thresholds", "CEP-style strategy engine"]),
    ("Frontend", BLUE, ["Next.js (App Router)", "React + TypeScript", "Tailwind CSS",
                        "Live SSE telemetry", "Recharts-style charts"]),
    ("Ops", RED, ["Docker + Compose", "systemd unit file", "Structured JSON logs",
                  "Secret redaction", "Mock + Webull adapters"]),
]
x0, cw, gx = 0.7, 2.95, 0.18
for i, (h_, c, items) in enumerate(groups):
    x = x0 + i * (cw + gx)
    rect(s, x, 2.1, cw, 4.4, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
    rect(s, x, 2.1, cw, 0.07, fill=c)
    text(s, x + 0.25, 2.35, cw - 0.4, 0.5, [[R(h_, 18, WHITE, True)]])
    paras = [[R("›  " + it, 13, MUTED, False, MONO)] for it in items]
    text(s, x + 0.25, 3.0, cw - 0.45, 3.3, paras, space_after=10)
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 11 — Live demo
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "Live Demo")
title(s, "Running it yourself")

# terminal panel
rect(s, 0.7, 2.05, 7.0, 4.5, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
text(s, 0.95, 2.2, 6, 0.35, [[R("● ● ●   terminal", 11, DIM, False, MONO)]])
cmds = [
    [R("# 1) Start the backend (paper mode)", 12.5, DIM, False, MONO)],
    [R("$ cd algo-trade", 13, EMERALD, False, MONO)],
    [R("$ python -m src.cli.main --mode paper", 13, EMERALD, False, MONO)],
    [R("", 6, DIM, False, MONO)],
    [R("# 2) Start the dashboard", 12.5, DIM, False, MONO)],
    [R("$ cd frontend && npm run dev", 13, EMERALD, False, MONO)],
    [R("  ▸ open http://localhost:3000", 12.5, BLUE, False, MONO)],
    [R("", 6, DIM, False, MONO)],
    [R("# 3) Run a backtest report", 12.5, DIM, False, MONO)],
    [R("$ python scripts/backtest.py \\", 13, EMERALD, False, MONO)],
    [R("      sample_data/minute_sample.csv", 13, EMERALD, False, MONO)],
]
text(s, 0.95, 2.75, 6.5, 3.6, cmds, space_after=6, line_spacing=1.0)

# right talking points
points = [
    ("No credentials needed", "Mock provider + mock broker run fully offline."),
    ("Watch the pipeline", "Logs show scan → model → signal → risk → fill."),
    ("Open the dashboard", "Live stats, chart, and circuit-breaker banner."),
    ("Show a backtest", "Print win rate & P&L over historical data."),
]
y = 2.1
for h_, b_ in points:
    rect(s, 8.0, y, 4.65, 1.0, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
    rect(s, 8.0, y, 0.07, 1.0, fill=EMERALD)
    text(s, 8.25, y + 0.13, 4.3, 0.4, [[R(h_, 14.5, WHITE, True)]])
    text(s, 8.25, y + 0.52, 4.3, 0.4, [[R(b_, 12, MUTED)]], line_spacing=1.0)
    y += 1.13
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 12 — Results / what the simulation produces
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=EMERALD)
kicker(s, "Results")
title(s, "What the simulation tells us")
text(s, 0.7, 1.95, 12, 0.6, [[R(
    "Every run produces measurable output — exactly what a Simulation & Modeling project "
    "needs to evaluate a model.", 15, MUTED)]], line_spacing=1.1)

metrics = [
    ("Signal count", "How often the model fired", EMERALD),
    ("Trades taken", "After risk + confirmation filters", BLUE),
    ("Win rate", "% of simulated trades profitable", AMBER),
    ("Net P&L", "Total simulated gain / loss", EMERALD),
]
x0, cw, gx = 0.7, 2.95, 0.18
for i, (h_, b_, c) in enumerate(metrics):
    x = x0 + i * (cw + gx)
    rect(s, x, 2.95, cw, 1.7, fill=PANEL, line=c, line_w=1.25, radius=True)
    text(s, x, 3.15, cw, 0.5, [[R(h_, 16, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, x, 3.7, cw, 0.7, [[R(b_, 12.5, MUTED)]], align=PP_ALIGN.CENTER, line_spacing=1.05)

text(s, 0.7, 5.1, 12, 1.4, [
    [R("Why this matters:", 17, EMERALD, True)],
    [R("Because the model is deterministic and the simulation is repeatable, we can "
       "change one parameter (say RSI threshold or position size), re-run, and directly "
       "compare results — the core experimental loop of modeling.", 15, MUTED)],
], line_spacing=1.15)
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 13 — Key takeaways
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s)
kicker(s, "Conclusion")
title(s, "Key takeaways")

takes = [
    ("Models a real-world system", "Market behavior captured as testable math (RSI/MACD/ATR).", EMERALD),
    ("Simulation-first design", "Backtest + paper modes let us evaluate safely, repeatably.", BLUE),
    ("Risk is part of the model", "Sizing & circuit breaker keep the simulation realistic.", RED),
    ("Clean, modular engineering", "Event-driven, decoupled, unit-tested components.", AMBER),
    ("Full-stack & demoable", "Python backend + live Next.js dashboard, runs offline.", BLUE),
]
y = 2.1
for h_, b_, c in takes:
    rect(s, 0.7, y, 11.92, 0.82, fill=PANEL, line=STROKE, line_w=1.0, radius=True)
    rect(s, 0.7, y, 0.1, 0.82, fill=c)
    chk = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), Inches(y + 0.24), Inches(0.34), Inches(0.34))
    _solid(chk, c); _noline(chk); chk.shadow.inherit = False
    text(s, 1.05, y + 0.2, 0.34, 0.34, [[R("✓", 14, BG, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.6, y + 0.12, 5.0, 0.6, [[R(h_, 16, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.5, y + 0.12, 6.0, 0.6, [[R(b_, 13.5, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.94
footer(s, nxt())


# ----------------------------------------------------------------------------
# Slide 14 — Thank you / Q&A
# ----------------------------------------------------------------------------
s = slide()
accent_bar(s, color=EMERALD)
rect(s, 8.4, 0, 4.93, 7.5, fill=PANEL)
rect(s, 8.4, 0, 0.02, 7.5, fill=STROKE)

text(s, 0.9, 2.6, 7.3, 2, [
    [R("Thank you", 60, WHITE, True)],
    [R("Questions & discussion", 24, EMERALD, True)],
])
text(s, 0.9, 4.7, 7.0, 1.2, [[R(
    "Algo-Trade — modeling the market, simulating the strategy, and keeping it safe. "
    "For educational & paper-trading use only.", 15, MUTED)]], line_spacing=1.15)
text(s, 0.9, 6.0, 8, 1.0, [
    [R("Presented by", 13, DIM, False, MONO)],
    [R("Muhammad Ahmad Saleem  F2022065202", 12.5, WHITE, False, MONO)],
    [R("Danish Salman  F2022065114", 12.5, WHITE, False, MONO)],
    [R("Irha Shoaib  F2021065136", 12.5, WHITE, False, MONO)],
], space_after=2, line_spacing=1.0)

text(s, 8.75, 2.7, 4.4, 4, [
    [R("RECAP", 13, EMERALD, True, MONO)],
    [R("", 8, MUTED)],
    [R("• Model:  RSI · MACD · ATR", 14, MUTED, False, MONO)],
    [R("• Sim:    backtest + paper", 14, MUTED, False, MONO)],
    [R("• Risk:   sizing + breaker", 14, MUTED, False, MONO)],
    [R("• Stack:  Python + Next.js", 14, MUTED, False, MONO)],
    [R("• Output: P&L · win rate", 14, MUTED, False, MONO)],
], space_after=8, line_spacing=1.05)
footer(s, nxt())


# ----------------------------------------------------------------------------
out_dir = Path(__file__).resolve().parent.parent / "docs"
out_dir.mkdir(exist_ok=True)
out_path = out_dir / "Algo-Trade-Simulation-and-Modeling.pptx"
try:
    prs.save(str(out_path))
except PermissionError:
    import time
    out_path = out_dir / f"Algo-Trade-Simulation-and-Modeling-{time.strftime('%H%M%S')}.pptx"
    prs.save(str(out_path))
    print("(original file was locked/open in PowerPoint — wrote a new copy)")
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
